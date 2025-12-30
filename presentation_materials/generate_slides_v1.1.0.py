import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def parse_markdown_slides(file_path):
    with open(file_path, 'r') as f:
        content = f.read()

    # Split by slide headers
    # Regex looks for "## SLIDE X.Y - Title"
    slide_pattern = re.compile(r'^## SLIDE [\d\.]+[a-z]? - (.+)$', re.MULTILINE)
    
    parts = slide_pattern.split(content)
    
    # parts[0] is intro text before first slide
    # parts[1] is Title 1, parts[2] is Content 1, parts[3] is Title 2...
    
    slides = []
    
    # Skip preamble (parts[0])
    # Iterate in pairs of (Title, Content)
    for i in range(1, len(parts), 2):
        title = parts[i].strip()
        body_text = parts[i+1].strip()
        slides.append({'title': title, 'content': body_text})
        
    return slides

def add_slide(prs, title, content):
    # Use blank layout for maximum control, or title and content
    # Layout 1 is usually Title + Content
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    
    # Set Title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    
    # Clean up content
    # Remove ``` markers
    clean_content = content.replace("```", "")
    
    # Remove "---" delimiters if present at start/end
    clean_content = clean_content.strip("- \n")

    # Set Content
    body_placeholder = slide.placeholders[1]
    text_frame = body_placeholder.text_frame
    text_frame.clear()  # Clear default empty paragraph
    
    lines = clean_content.split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        p = text_frame.add_paragraph()
        
        # Check for bullet points
        if line.startswith('•') or line.startswith('- ') or line.startswith('* '):
            p.text = line.lstrip('•-* ')
            p.level = 0 
        else:
            p.text = line
            p.level = 0
            
        # Basic font formatting
        p.font.size = Pt(18)
        if "━━━━━━━━" in line:
            p.text = "" # Skip visual separators

def create_presentation(input_md, output_pptx):
    prs = Presentation()
    
    slides_data = parse_markdown_slides(input_md)
    
    print(f"Found {len(slides_data)} slides. Generating...")
    
    for slide_data in slides_data:
        add_slide(prs, slide_data['title'], slide_data['content'])
        
    prs.save(output_pptx)
    print(f"Saved presentation to {output_pptx}")

if __name__ == "__main__":
    input_file = "Universal-Life-Science-and-Clinical-Skills-/presentation_materials/PRESENTATION_SLIDES.md"
    # Version 1.1.0
    output_file = "Universal-Life-Science-and-Clinical-Skills-/presentation_materials/Universal_Biomedical_Skills_Presentation_v1.1.0.pptx"
    
    create_presentation(input_file, output_file)
